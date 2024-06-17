from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

class PPTXBuilder:
    def __init__(self, output_filename='output_presentation.pptx'):
        self.prs = Presentation()
        self.output_filename = output_filename

    def add_slide_with_image(self, img_path, left=Inches(1), top=Inches(1), height=Inches(5)):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        with Image.open(img_path) as img:
            img_width, img_height = img.size
        aspect_ratio = img_width / img_height
        width = height * aspect_ratio
        slide.shapes.add_picture(img_path, left, top, width, height)
        return slide

    def add_text_to_slide(self, slide, text, left=Inches(1), top=Inches(6), width=Inches(8), height=Inches(1)):
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Inches(0.5)  # Adjust the font size as needed

    def save(self):
        self.prs.save(self.output_filename)
        print(f"Presentation saved as {self.output_filename}")

def create_presentation(main_folder_path, epc_folder_name, gsa_folder_name, output_filename='output_presentation.pptx'):
    # Create a PPTXBuilder object
    ppt_builder = PPTXBuilder(output_filename)

    # Check for EPC folder and add slides
    epc_folder_path = os.path.join(main_folder_path, epc_folder_name)
    if os.path.isdir(epc_folder_path):
        for filename in os.listdir(epc_folder_path):
            if filename.endswith('.png'):
                img_path = os.path.join(epc_folder_path, filename)
                slide = ppt_builder.add_slide_with_image(img_path)
                ppt_builder.add_text_to_slide(slide, f"Image from EPC: {filename}")

    # Check for GSA folder and add slides
    gsa_folder_path = os.path.join(main_folder_path, gsa_folder_name)
    if os.path.isdir(gsa_folder_path):
        for filename in os.listdir(gsa_folder_path):
            if filename.endswith('.png'):
                img_path = os.path.join(gsa_folder_path, filename)
                slide = ppt_builder.add_slide_with_image(img_path)
                ppt_builder.add_text_to_slide(slide, f"Image from GSA: {filename}")

    # Save the presentation
    ppt_builder.save()

# Define the main folder and subfolders
main_folder_path = 'result_2023'
epc_folder_name = 'EPC'
gsa_folder_name = 'GSA'

# Create the PowerPoint presentation
create_presentation(main_folder_path, epc_folder_name, gsa_folder_name)
